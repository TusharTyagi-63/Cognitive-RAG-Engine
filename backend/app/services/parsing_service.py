"""
backend/app/services/parsing_service.py
=======================================
Extracts plain text from raw files stored on disk.

Supports: PDF, DOCX, PPTX, XLSX, CSV, HTML, XML, JSON, RTF,
          Images (via Gemini Vision AI), Markdown, plain text,
          and all common source code / config formats.
"""
import csv
import json
from pathlib import Path

import pymupdf

from backend.app.utils.exceptions import BadRequestException


# Extensions that can be read as plain text (source code, configs, logs)
_TEXT_EXTENSIONS = {
    ".txt", ".md", ".py", ".js", ".ts", ".java", ".c", ".cpp", ".go",
    ".rs", ".rb", ".php", ".sh", ".sql", ".yaml", ".yml", ".toml",
    ".ini", ".cfg", ".log",
}

# Image extensions — parsed via Gemini Vision AI
_IMAGE_EXTENSIONS = {
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif", ".svg",
}


class ParsingService:
    @staticmethod
    def extract_text(file_path: Path, content_type: str) -> str:
        """
        Routes the file to the correct parser based on MIME type or extension.
        Returns the extracted plain text as a single string.
        """
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = file_path.suffix.lower()

        try:
            # ── PDF ──────────────────────────────────────────────────────
            if content_type == "application/pdf" or ext == ".pdf":
                return ParsingService._parse_pdf(file_path)

            # ── Microsoft Office ─────────────────────────────────────────
            if ext == ".docx" or content_type in (
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ):
                return ParsingService._parse_docx(file_path)

            if ext == ".pptx" or content_type in (
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ):
                return ParsingService._parse_pptx(file_path)

            if ext == ".xlsx" or content_type in (
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ):
                return ParsingService._parse_xlsx(file_path)

            # ── CSV ──────────────────────────────────────────────────────
            if content_type in ("text/csv", "application/csv") or ext == ".csv":
                return ParsingService._parse_csv(file_path)

            # ── HTML ─────────────────────────────────────────────────────
            if ext in (".html", ".htm") or content_type == "text/html":
                return ParsingService._parse_html(file_path)

            # ── XML ──────────────────────────────────────────────────────
            if ext == ".xml" or content_type in ("application/xml", "text/xml"):
                return ParsingService._parse_xml(file_path)

            # ── JSON ─────────────────────────────────────────────────────
            if ext == ".json" or content_type == "application/json":
                return ParsingService._parse_json(file_path)

            # ── RTF ──────────────────────────────────────────────────────
            if ext == ".rtf" or content_type == "application/rtf":
                return ParsingService._parse_rtf(file_path)

            # ── Images (Gemini Vision AI) ────────────────────────────────
            if ext in _IMAGE_EXTENSIONS or content_type.startswith("image/"):
                return ParsingService._parse_image(file_path)

            # ── Plain text / source code / config ────────────────────────
            if ext in _TEXT_EXTENSIONS or content_type.startswith("text/"):
                return ParsingService._parse_text(file_path)

            # ── Fallback: attempt plain text ─────────────────────────────
            return ParsingService._parse_text(file_path)

        except BadRequestException:
            raise
        except Exception as e:
            raise BadRequestException(f"Failed to parse file: {str(e)}")

    # ──────────────────────────────────────────────────────────────────────
    # Individual Parsers
    # ──────────────────────────────────────────────────────────────────────

    @staticmethod
    def _parse_pdf(file_path: Path) -> str:
        """Extracts text from a PDF file using PyMuPDF."""
        text_blocks = []
        with pymupdf.open(file_path) as doc:
            for page in doc:
                text_blocks.append(page.get_text())
        return "\n\n".join(text_blocks)

    @staticmethod
    def _parse_docx(file_path: Path) -> str:
        """Extracts text from a Word (.docx) file including paragraphs and tables."""
        from docx import Document

        doc = Document(str(file_path))
        parts = []

        # Extract paragraphs
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)

        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))

        return "\n\n".join(parts)

    @staticmethod
    def _parse_pptx(file_path: Path) -> str:
        """Extracts text from a PowerPoint (.pptx) file, slide by slide."""
        from pptx import Presentation

        prs = Presentation(str(file_path))
        parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)

                # Extract text from tables inside slides
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_texts.append(" | ".join(cells))

            if slide_texts:
                parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))

        return "\n\n".join(parts)

    @staticmethod
    def _parse_xlsx(file_path: Path) -> str:
        """Extracts text from an Excel (.xlsx) file, sheet by sheet."""
        from openpyxl import load_workbook

        wb = load_workbook(str(file_path), read_only=True, data_only=True)
        parts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_lines = [f"--- Sheet: {sheet_name} ---"]

            for row in ws.iter_rows(values_only=True):
                cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if cells:
                    sheet_lines.append(" | ".join(cells))

            if len(sheet_lines) > 1:  # More than just the header
                parts.append("\n".join(sheet_lines))

        wb.close()
        return "\n\n".join(parts)

    @staticmethod
    def _parse_csv(file_path: Path) -> str:
        """Extracts text from a CSV file, formatting it cleanly."""
        lines = []
        with open(file_path, mode='r', encoding='utf-8', errors='replace') as csvfile:
            reader = csv.reader(csvfile)
            for row in reader:
                # Join non-empty cells with a tab or comma for readable text
                lines.append(" | ".join([cell.strip() for cell in row if cell.strip()]))
        return "\n".join(lines)

    @staticmethod
    def _parse_html(file_path: Path) -> str:
        """Extracts readable text from an HTML file, stripping tags."""
        from bs4 import BeautifulSoup

        with open(file_path, mode='r', encoding='utf-8', errors='replace') as f:
            soup = BeautifulSoup(f.read(), 'lxml')

        # Remove script and style elements
        for tag in soup(['script', 'style', 'nav', 'footer', 'header']):
            tag.decompose()

        return soup.get_text(separator='\n', strip=True)

    @staticmethod
    def _parse_xml(file_path: Path) -> str:
        """Extracts text content from an XML file."""
        from bs4 import BeautifulSoup

        with open(file_path, mode='r', encoding='utf-8', errors='replace') as f:
            soup = BeautifulSoup(f.read(), 'lxml-xml')

        return soup.get_text(separator='\n', strip=True)

    @staticmethod
    def _parse_json(file_path: Path) -> str:
        """Pretty-prints a JSON file for AI readability."""
        with open(file_path, mode='r', encoding='utf-8', errors='replace') as f:
            data = json.load(f)
        return json.dumps(data, indent=2, ensure_ascii=False)

    @staticmethod
    def _parse_rtf(file_path: Path) -> str:
        """Extracts plain text from an RTF file."""
        from striprtf.striprtf import rtf_to_text

        with open(file_path, mode='r', encoding='utf-8', errors='replace') as f:
            rtf_content = f.read()
        return rtf_to_text(rtf_content)

    @staticmethod
    def _parse_text(file_path: Path) -> str:
        """Reads plain text, markdown, source code, or config files."""
        with open(file_path, mode='r', encoding='utf-8', errors='replace') as f:
            return f.read()

    @staticmethod
    def _parse_image(file_path: Path) -> str:
        """
        Extracts text and descriptions from an image using Google Gemini Vision.
        This handles text extraction (OCR), diagram reading, chart analysis,
        and general image description — far superior to traditional OCR.
        """
        import google.generativeai as genai
        from PIL import Image
        from backend.app.core.config import settings

        if not settings.GEMINI_API_KEY:
            raise BadRequestException("GEMINI_API_KEY is not configured. Cannot process images.")

        genai.configure(api_key=settings.GEMINI_API_KEY)
        model = genai.GenerativeModel('gemini-2.0-flash')

        img = Image.open(file_path)

        prompt = (
            "Analyze this image thoroughly. Extract ALL text visible in the image exactly as written. "
            "If the image contains diagrams, charts, tables, or figures, describe them in detail "
            "including all data points, labels, and relationships. "
            "If it's a screenshot of code, extract the code exactly. "
            "If it's a photograph, describe what you see in detail. "
            "Format the output as clean, readable text."
        )

        response = model.generate_content([prompt, img])
        return response.text
